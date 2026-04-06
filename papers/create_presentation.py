"""
create_presentation.py

Generates a .pptx presentation:
  "Qualitative Reasoning: History and the Rise of LLMs"

Audience: CS professors unfamiliar with QR
Slides:   ~40 content slides + references
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── colour palette ────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x37, 0x5E)   # title backgrounds
MID_BLUE    = RGBColor(0x2E, 0x5E, 0x9B)   # section dividers
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF7)   # body backgrounds
ACCENT      = RGBColor(0xE8, 0x6B, 0x00)   # highlights / callout
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── low-level helpers ─────────────────────────────────────────────────────────

def fill_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide, DARK_BLUE)

    # large colour bar
    add_rect(slide, Inches(0), Inches(2.2), SLIDE_W, Inches(3.1), MID_BLUE)

    add_textbox(slide,
                Inches(0.6), Inches(2.35), Inches(12.1), Inches(2.0),
                title, font_size=40, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    if subtitle:
        add_textbox(slide,
                    Inches(0.6), Inches(4.4), Inches(12.1), Inches(0.8),
                    subtitle, font_size=22, bold=False, color=LIGHT_BLUE,
                    align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5),
                "Qualitative Reasoning Workshop (QR) — 2025",
                font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    return slide


def add_section_divider(title, subtitle=""):
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide, MID_BLUE)

    add_rect(slide, Inches(0), Inches(0), Inches(0.3), SLIDE_H, ACCENT)

    add_textbox(slide,
                Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.4),
                title, font_size=36, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide,
                    Inches(0.8), Inches(4.1), Inches(11.7), Inches(1.0),
                    subtitle, font_size=20, color=LIGHT_BLUE,
                    align=PP_ALIGN.LEFT)
    return slide


def add_content_slide(title, bullets, notes=""):
    """Standard two-zone slide: dark header + light body with bullet list."""
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide, LIGHT_GRAY)

    # header bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), DARK_BLUE)
    add_textbox(slide,
                Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.95),
                title, font_size=26, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)

    # body
    body_top  = Inches(1.25)
    body_left = Inches(0.4)
    body_w    = Inches(12.5)
    body_h    = Inches(5.9)

    txBox = slide.shapes.add_textbox(body_left, body_top, body_w, body_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.level = level
        p.space_before = Pt(4 if level == 0 else 2)
        indent = "    " * level
        bullet_char = "•" if level == 0 else "–"
        run = p.add_run()
        run.text = f"{indent}{bullet_char}  {text}"
        run.font.size = Pt(17 if level == 0 else 15)
        run.font.bold = (level == 0)
        run.font.color.rgb = DARK_BLUE if level == 0 else BLACK

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_two_column_slide(title, left_header, left_items,
                         right_header, right_items):
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide, LIGHT_GRAY)

    # header
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), DARK_BLUE)
    add_textbox(slide,
                Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.95),
                title, font_size=26, bold=True, color=WHITE)

    col_top = Inches(1.3)
    col_h   = Inches(5.8)

    for col_left, header, items in [
        (Inches(0.3),  left_header,  left_items),
        (Inches(6.85), right_header, right_items),
    ]:
        col_w = Inches(6.1)
        # column header
        add_rect(slide, col_left, col_top, col_w, Inches(0.42), MID_BLUE)
        add_textbox(slide, col_left + Inches(0.1),
                    col_top + Inches(0.04),
                    col_w - Inches(0.2), Inches(0.38),
                    header, font_size=16, bold=True, color=WHITE)

        # column body
        txBox = slide.shapes.add_textbox(
            col_left + Inches(0.1),
            col_top + Inches(0.52),
            col_w - Inches(0.2),
            col_h - Inches(0.52)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = f"•  {item}"
            run.font.size = Pt(15)
            run.font.color.rgb = BLACK

    return slide


def add_code_slide(title, code_text, caption=""):
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide, LIGHT_GRAY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), DARK_BLUE)
    add_textbox(slide,
                Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.95),
                title, font_size=26, bold=True, color=WHITE)

    # dark code box
    code_top = Inches(1.25)
    code_h   = Inches(5.0)
    add_rect(slide, Inches(0.4), code_top, Inches(12.5), code_h,
             RGBColor(0x1E, 0x1E, 0x1E))
    add_textbox(slide,
                Inches(0.55), code_top + Inches(0.1),
                Inches(12.2), code_h - Inches(0.2),
                code_text, font_size=14, color=RGBColor(0xD4, 0xD4, 0xD4),
                bold=False, italic=False)

    if caption:
        add_textbox(slide,
                    Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.6),
                    caption, font_size=14, italic=True, color=RGBColor(0x55, 0x55, 0x55))
    return slide


def add_quote_slide(title, quote, attribution=""):
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide, DARK_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), MID_BLUE)
    add_textbox(slide,
                Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.95),
                title, font_size=26, bold=True, color=WHITE)

    # large quote mark decoration
    add_textbox(slide,
                Inches(0.3), Inches(1.1), Inches(1.5), Inches(2.0),
                "\u201c", font_size=120, bold=True,
                color=RGBColor(0x4A, 0x7A, 0xC0))

    add_textbox(slide,
                Inches(1.2), Inches(1.8), Inches(10.8), Inches(4.0),
                quote, font_size=20, color=WHITE, italic=True,
                align=PP_ALIGN.LEFT)

    if attribution:
        add_textbox(slide,
                    Inches(1.2), Inches(5.9), Inches(10.8), Inches(0.7),
                    f"— {attribution}", font_size=16,
                    color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE CONTENT
# ═══════════════════════════════════════════════════════════════════════════════

# ── 1. Title ──────────────────────────────────────────────────────────────────
add_title_slide(
    "Qualitative Reasoning:\nHistory and the Rise of LLMs",
    "From Sign Algebra to GPT-4 — 60 Years of Commonsense Physics"
)

# ── 2. Agenda ─────────────────────────────────────────────────────────────────
add_content_slide("Agenda", [
    "Part 1 — What is Qualitative Reasoning?",
    ("Motivation: reasoning without equations", 1),
    ("Core representations and formalisms", 1),
    ("Key algorithms: QSIM, QP Theory, compositional modeling", 1),
    ("Real-world applications", 1),
    "Part 2 — Pre-AI Origins (1960s–1970s)",
    ("Sign algebra in economics; control theory", 1),
    ("Hayes' Naïve Physics Manifesto", 1),
    "Part 3 — The AI Era (1977–present)",
    ("Foundational 1984 papers; the QR workshop series", 1),
    ("Spatial reasoning, education tools, QR software", 1),
    "Part 4 — LLMs Meet QR (2023–2025)",
    ("Benchmarking LLMs on spatial tasks", 1),
    ("Hybrid neuro-symbolic architectures", 1),
    ("The road ahead", 1),
    "References",
])

# ══════════════════════════════════════════════════════════════════════════════
# PART 1 — WHAT IS QUALITATIVE REASONING?
# ══════════════════════════════════════════════════════════════════════════════

add_section_divider(
    "Part 1 — What Is Qualitative Reasoning?",
    "Reasoning about continuous phenomena without full equations"
)

# ── 3. Motivation ─────────────────────────────────────────────────────────────
add_content_slide("The Central Motivation", [
    "Two empirical observations drive the field:",
    ("People routinely draw useful conclusions without numerical equations",     1),
    ("Scientists and engineers use qualitative reasoning when setting up problems and interpreting results", 1),
    "Example: without any numbers, you can conclude…",
    ("Heating a gas in a fixed container raises its pressure", 1),
    ("Releasing a ball near a slope makes it roll downward", 1),
    ("Removing a load from a spring shortens it", 1),
    "Goal: give computers the same commonsense reasoning ability",
    ("Formal, symbolic — not statistical or learned from data", 1),
    ("Produce causal, explainable predictions", 1),
    ("Work with incomplete information (no exact values required)", 1),
])

# ── 4. What QR Is NOT ─────────────────────────────────────────────────────────
add_content_slide("What QR Is and Is Not", [
    "QR IS:",
    ("Symbolic AI: uses formal representations of qualitative states", 1),
    ("Causal: models mechanisms, not just correlations", 1),
    ("Sound under incomplete information: works when only + / 0 / − is known", 1),
    ("Explainable: every conclusion traces to model structure", 1),
    "QR IS NOT:",
    ("Fuzzy logic (no membership functions)", 1),
    ("Numerical simulation (no ODE solvers)", 1),
    ("Machine learning (no training data required)", 1),
    ("Approximate or heuristic — it is formally rigorous", 1),
    "Closest relatives in CS: model-based reasoning, constraint programming, knowledge representation",
])

# ── 5. The Representation Hierarchy ───────────────────────────────────────────
add_content_slide("QR Representations: A Hierarchy of Precision", [
    "Sign Algebra (roughest)  →  Quantity Spaces  →  Numerical Values (finest)",
    "Sign algebra: every quantity is + (positive), 0 (zero), or − (negative)",
    ("Add: (+) + (+) = (+),   (+) + (−) = ?,   (+) + (0) = (+)", 1),
    ("Multiply: (+) × (−) = (−),   (−) × (−) = (+)", 1),
    ("Used for comparative statics: if X increases, what happens to Y?", 1),
    "Quantity spaces: ordered sets of landmark values",
    ("e.g., {0, boiling, above_boiling} for water temperature", 1),
    ("Captures threshold effects and qualitative transitions", 1),
    ("Magnitudes and derivatives each take values in a quantity space", 1),
    "Monotonic function constraints: Y = M+(X) means Y increases when X increases",
    ("No equation needed — just the direction of influence", 1),
])

# ── 6. Key Formalisms ─────────────────────────────────────────────────────────
add_two_column_slide(
    "Core Formalisms in QR",
    "Confluences (de Kleer & Brown, 1984)",
    [
        "Differential equations over sign algebra",
        "Written as: δP + δV = 0 (Boyle's law, qualitatively)",
        "δ means 'direction of change'",
        "Used for device modeling (circuits, pipes)",
        "Device ontology: components + connections",
        "Envisionment: all reachable qualitative states",
    ],
    "Qualitative Differential Equations (Kuipers, 1986)",
    [
        "QDE: set of constraints over quantities and their derivatives",
        "QSIM operators: ADD, MULT, DERIV, M+, M−",
        "Simulation traces a tree of possible behaviors",
        "Handles ambiguity — multiple consistent futures",
        "QSIM algorithm: globally consistent state propagation",
        "Behavior prediction without any numerical values",
    ]
)

# ── 7. Qualitative Process Theory ─────────────────────────────────────────────
add_content_slide("Qualitative Process Theory (Forbus, 1984)", [
    "Central insight: continuous change is caused by processes",
    "A process model fragment specifies:",
    ("Individuals: objects participating in the process", 1),
    ("Preconditions: qualitative conditions for the process to exist", 1),
    ("Quantity conditions: inequalities on quantity magnitudes", 1),
    ("Relations: direct influences (I+, I−) and proportionalities (P+, P−)", 1),
    "Direct influences set the sign of a quantity's derivative",
    ("Heat flow into an object → I+(Temperature, Heat-Flow)", 1),
    "Proportionalities propagate directional change through the system",
    ("Higher temperature → proportionally higher pressure (P+)", 1),
    "Process ontology vs. device ontology (de Kleer):",
    ("QPT: causality flows from active processes, not component structure", 1),
    ("Both remain important — hybrid approaches exist", 1),
])

# ── 8. QSIM ───────────────────────────────────────────────────────────────────
add_content_slide("QSIM: Qualitative Simulation Algorithm (Kuipers)", [
    "Input: a Qualitative Differential Equation (QDE) + initial state",
    "Output: a tree (or graph) of all qualitatively consistent behaviors",
    "Algorithm sketch:",
    ("1. Represent each quantity as (magnitude, derivative) in its quantity space", 1),
    ("2. Apply global filters: P+/P− consistency, ADD/MULT constraints, DERIV sign", 1),
    ("3. Generate all successor states — transitions happen at landmark crossings", 1),
    ("4. Prune states violating any constraint; recurse on survivors", 1),
    "Key properties:",
    ("Soundness: every real behavior is contained in the QSIM output", 1),
    ("Completeness: QSIM may predict spurious behaviors (false positives), never misses real ones", 1),
    ("Ambiguity is a feature: captures genuine uncertainty in model", 1),
    "Extension: Order-of-Magnitude reasoning (Mavrovouniotis & Stephanopoulos, 1988)",
    ("Ranks quantities: negligible, comparable, dominant — finer than sign algebra", 1),
])

# ── 9. Compositional Modeling ─────────────────────────────────────────────────
add_content_slide("Compositional Modeling (Falkenhainer & Forbus, 1991)", [
    "Problem: hand-crafting a full model for each system is expensive and error-prone",
    "Solution: build models from a domain theory — a library of model fragments",
    "Model fragments encode reusable physical knowledge:",
    ("Static model fragments: define entities and their properties", 1),
    ("Process model fragments: define causal processes (e.g., liquid-flow, heat-transfer)", 1),
    ("Agent model fragments: define external interventions", 1),
    "Modeling assumptions control which fragments are relevant:",
    ("Ontological assumptions: what kinds of things exist?", 1),
    ("Perspective assumptions: which aspects of behavior matter?", 1),
    ("Granularity assumptions: at what scale is the model operating?", 1),
    "Result: given a scenario description, the system assembles the model automatically",
    ("Supports reuse across domains: same heat-flow fragment applies to ovens, engines, ecosystems", 1),
])

# ── 10. Qualitative Spatial Reasoning ─────────────────────────────────────────
add_content_slide("Qualitative Spatial Reasoning and RCC-8", [
    "Not all QR is about quantities — space and shape matter too",
    "RCC-8 (Region Connection Calculus): 8 qualitative relations between spatial regions",
    ("DC: Disconnected  |  EC: Externally Connected  |  PO: Partially Overlapping", 1),
    ("TPP: Tangential Proper Part  |  NTPP: Non-Tangential Proper Part", 1),
    ("TPPi, NTPPi (inverses)  |  EQ: Equal", 1),
    "Composition table: if A EC B and B DC C, what is the relation between A and C?",
    ("27 base cases; each cell is a set of possible relations", 1),
    ("Formal reasoning: sound and complete (unlike LLMs, as we will see)", 1),
    "Conceptual neighborhood: which relations are topologically adjacent?",
    ("DC can transition to EC smoothly; DC cannot jump directly to NTPP", 1),
    "Applications: robotics, geographic information systems (GIS), CAD, visual reasoning",
    "CogSketch (Forbus et al.): diagrammatic reasoning using qualitative spatial representations",
])

# ══════════════════════════════════════════════════════════════════════════════
# PART 2 — PRE-AI ORIGINS
# ══════════════════════════════════════════════════════════════════════════════

add_section_divider(
    "Part 2 — Pre-AI Origins (1960s–1970s)",
    "Sign algebra, comparative statics, and the Naïve Physics manifesto"
)

# ── 11. Economics: Sign Algebra ────────────────────────────────────────────────
add_content_slide("Pre-AI Origins: Sign Algebra in Economics", [
    "The oldest form of QR traces to mathematical economics, not AI",
    "Lancaster (1962) — 'The Solution of Qualitative Comparative Static Problems'",
    ("Introduced formal sign algebra for economic models", 1),
    ("Comparative statics: if a parameter changes, which way does equilibrium shift?", 1),
    ("Signed Jacobian matrices: each entry is +, −, or 0", 1),
    "Samuelson (1947) — Foundations of Economic Analysis",
    ("Qualitative restrictions as the basis for testable hypotheses", 1),
    ("Foundation for what became 'qualitative economics'", 1),
    "Gorman (1964), Bassett et al. (1968): extended sign algebra to dynamic models",
    "Key limitation: static models only — no concept of time or state transitions",
    "AI would later add the dynamic dimension, causal mechanisms, and explanation",
])

# ── 12. Control Theory ────────────────────────────────────────────────────────
add_content_slide("Pre-AI Origins: Control Theory (1960s–1970s)", [
    "Signed-graph / influence-diagram analysis in control engineering",
    "Loop analysis (Puccia & Levins, 1985 — consolidating prior work):",
    ("Represent systems as signed directed graphs", 1),
    ("Edge sign indicates direction of influence: + or −", 1),
    ("Determine system-level stability from graph structure", 1),
    "Quirk, Ruppert & Saposnik (1968) — 'Qualitative economics and the stability of equilibrium'",
    ("Conditions on sign patterns sufficient for unique equilibrium", 1),
    "Roberts (1971, 1976) — pulse processes on signed digraphs",
    ("Models ecological and social systems qualitatively", 1),
    ("First formal treatment of qualitative dynamics", 1),
    "Connection to AI: these formalisms had no inference engine, no explanation capability",
    ("AI contributions: formal soundness proofs, simulation algorithms, knowledge representation", 1),
])

# ── 13. Hayes' Manifesto ──────────────────────────────────────────────────────
add_quote_slide(
    "The Naïve Physics Manifesto — Patrick Hayes (1977)",
    "We need to formalise a sizable portion of commonsense knowledge about the physical world: "
    "about objects, their shapes, space, time, movement, substances, and so on. "
    "The naive physicist's world is rich, complex and various. It will take a large and careful effort "
    "to formalise it. But it can and should be done.",
    "Patrick Hayes, 'In Defence of Logic', IJCAI 1977"
)

# ── 14. Hayes continued ───────────────────────────────────────────────────────
add_content_slide("Naïve Physics: The Research Agenda", [
    "Hayes (1977, 1979, 1985) articulated what AI needed to do:",
    ("Model the everyday physical world in formal logic", 1),
    ("Cover: solid objects, liquids, shapes, space, time, movement, change", 1),
    ("Make implicit commonsense knowledge explicit and computable", 1),
    "Hayes' Histories (1985): formal theory of liquid behavior",
    ("Spatial-temporal histories: regions of space-time occupied by matter", 1),
    ("Portals: points of connection through which fluid flows", 1),
    ("'Naïve physics is hard' — history theory alone ran to dozens of axioms", 1),
    "Impact on AI:",
    ("Catalyzed the QR research program at MIT (de Kleer, Forbus)", 1),
    ("Established that commonsense physical reasoning is a tractable research target", 1),
    ("Framed the problem as knowledge engineering, not statistical learning", 1),
    "Still relevant today: the 'grounding problem' for LLMs echoes Hayes' original concern",
])

# ══════════════════════════════════════════════════════════════════════════════
# PART 3 — THE AI ERA
# ══════════════════════════════════════════════════════════════════════════════

add_section_divider(
    "Part 3 — The AI Era (1977–present)",
    "Foundational papers, algorithms, and real-world applications"
)

# ── 15. The Founding Papers ───────────────────────────────────────────────────
add_content_slide("The Founding Moment: 1984 AI Journal Special Issue", [
    "Three simultaneous landmark papers, all in Artificial Intelligence vol. 24 (1984):",
    "de Kleer & Brown — 'A Qualitative Physics Based on Confluences'",
    ("Confluences: algebraic constraints over sign algebra", 1),
    ("Device-centered: model systems as networks of interacting components", 1),
    ("Envisionment: exhaustively enumerate all qualitative states", 1),
    "Forbus — 'Qualitative Process Theory'",
    ("Process-centered: continuous change is caused by active processes", 1),
    ("I+, I−, P+, P− relations; formal condition system for process activity", 1),
    "Kuipers — 'Commonsense Reasoning about Causality: Deriving Behavior from Structure'",
    ("Earliest version of QSIM ideas — behavior from structural description", 1),
    ("Extended later into full QSIM (1986) with QDEs", 1),
    "These papers defined the field. QR has held annual workshops since 1987.",
])

# ── 16. The QR Workshop Series ────────────────────────────────────────────────
add_content_slide("The QR Workshop Series (1987–present)", [
    "Annual workshop co-located with major AI conferences (IJCAI, AAAI, ECAI)",
    "Now in its 39th year — one of AI's longest-running specialized workshops",
    "Workshop scope has expanded from physics to:",
    ("Engineering design and diagnosis", 1),
    ("Education and intelligent tutoring", 1),
    ("Ecological and biological systems", 1),
    ("Social and economic systems", 1),
    ("Autonomous vehicles and robotics", 1),
    ("Natural language understanding", 1),
    ("LLM evaluation and integration (most recently)", 1),
    "Key software artifacts developed by the community:",
    ("QSIM (Kuipers, UT Austin)", 1),
    ("QPE / Garp3 / DynaLearn (Bredeweg et al., University of Amsterdam)", 1),
    ("CogSketch (Forbus, Northwestern)", 1),
    ("NextKB / OpenCyc (Forbus, Northwestern)", 1),
])

# ── 17. Key Algorithms and Their Relations ─────────────────────────────────────
add_content_slide("Key QR Algorithms at a Glance", [
    "Sign algebra propagation (1960s–1980s)",
    ("Simple matrix operations; no temporal dimension", 1),
    "QSIM (Kuipers 1986–2001)",
    ("Globally consistent state generation; soundness guarantee; handles ambiguity", 1),
    ("Extended: Q2 (Kuipers et al.) — reduces spurious predictions", 1),
    "QPT simulation (Forbus 1984–present)",
    ("Process activation and deactivation; episodes with landmark crossings", 1),
    "Compositional modeling / QSIM+ (Falkenhainer & Forbus 1991)",
    ("Model assembly from domain theory; modeling assumptions", 1),
    "Qualitative causal modeling (de Kleer 1986)",
    ("Causal ordering; device-level diagnosis (GDE algorithm)", 1),
    "Order-of-magnitude reasoning (Mavrovouniotis 1988)",
    ("Five levels: negligible ≪ comparable ≈ dominant; finer than sign algebra", 1),
    "Constraint propagation in QR (Leler 1988; Dague et al. 1992)",
    ("Arc-consistency for qualitative constraint networks", 1),
])

# ── 18. Real-World Applications (1990s–2000s) ─────────────────────────────────
add_content_slide("Real-World Applications: 1990s–2000s", [
    "Commercial deployment began in the early 1990s:",
    "Xerox DC-6090 photocopier (Hamscher, ~1990)",
    ("QR-based fault isolation from symptom observations", 1),
    ("First large-scale commercial use of qualitative diagnosis", 1),
    "Automotive FMEA — Mentor Graphics / Sherlock (~1995)",
    ("Failure Mode and Effects Analysis for circuit boards", 1),
    ("Qualitative causal propagation of failure effects through designs", 1),
    "NASA Deep Space One — Livingstone system (Williams & Nayak, 1996)",
    ("Model-based autonomous reconfiguration after fault detection", 1),
    ("Deployed on actual spacecraft: first AI-based autonomous spacecraft", 1),
    "Process control and supervision (chemical plants, power systems)",
    ("COGSIM, FALCON: qualitative models for supervisory control", 1),
    "Water treatment decision support (Struss, QR 2023)",
    ("QP model fragments for coagulation, filtration, sedimentation", 1),
])

# ── 19. Applications: Education ───────────────────────────────────────────────
add_content_slide("Applications: Education and Cognitive Science", [
    "DynaLearn / Garp3 (Bredeweg, University of Amsterdam — 2000s–present)",
    ("Interactive qualitative modeling environment for students", 1),
    ("Students construct causal models, run simulations, explore scenarios", 1),
    ("Used for stellar physics, orbital motion, thermoregulation, photoelectric effect, global warming", 1),
    ("Deployed in secondary and higher education across Europe", 1),
    "CogSketch (Forbus, Northwestern)",
    ("Qualitative spatial reasoning over student-drawn diagrams", 1),
    ("Tutors physics problem-solving through diagrammatic QR", 1),
    "Why QR for education?",
    ("Causal models align with how students should think, not just numerical answers", 1),
    ("Incorrect model fragments reveal student misconceptions explicitly", 1),
    ("Model construction is itself the learning activity", 1),
    "GarpN (Kragten et al., QR 2024): bridges qualitative and quantitative modeling",
    ("Converts DynaLearn models to numerical simulations for comparison", 1),
])

# ── 20. The 2006 Visions ──────────────────────────────────────────────────────
add_content_slide("20-Year Visions from QR 2006 (Price et al.)", [
    "In 2006, the QR community articulated six grand challenge visions:",
    "Science-bot: autonomous scientific discovery using QR for hypothesis generation",
    "Virtual vehicle: comprehensive qualitative model of an entire automobile system",
    "Complex natural systems: QR models of ecological and climate systems at multiple scales",
    "4D medical monitoring: qualitative temporal abstraction of continuous physiological signals",
    "Autonomous problem solvers: AI that formulates, solves, and explains physical problems",
    "Encyclopedia of human mental models: QR formalization of naïve physics knowledge",
    "Technology gaps identified in 2006 still relevant today:",
    ("Hybrid modeling: integrating qualitative and quantitative methods seamlessly", 1),
    ("Multi-level modeling: coherent models across scales of abstraction", 1),
    ("Automated model generation: reducing the knowledge engineering bottleneck", 1),
    "Reflection: LLMs now address part of 'automated model generation' — but introduce new challenges",
])

# ══════════════════════════════════════════════════════════════════════════════
# PART 4 — LLMs IN QR
# ══════════════════════════════════════════════════════════════════════════════

add_section_divider(
    "Part 4 — LLMs Meet Qualitative Reasoning (2023–2025)",
    "Benchmarks, hybrid architectures, and open questions"
)

# ── 21. The LLM Moment for QR ─────────────────────────────────────────────────
add_content_slide("The LLM Revolution Reaches QR", [
    "GPT-4 (2023) performed impressively on many reasoning benchmarks — but what about QR?",
    "The QR community started asking three questions:",
    ("Can LLMs perform formal qualitative spatial / causal reasoning?", 1),
    ("Can LLMs assist in building qualitative models (reducing KE bottleneck)?", 1),
    ("Should LLMs replace formal QR representations entirely?", 1),
    "Three years of QR workshops (2023, 2024, 2025) provide a systematic answer",
    "Summary verdict so far:",
    ("LLMs are useful front-ends (natural language → structure)", 1),
    ("LLMs fail at formal QR reasoning without symbolic support", 1),
    ("Large Reasoning Models (LRMs: o1, o3) dramatically close the gap on benchmarks", 1),
    ("Hybrid neuro-symbolic architectures consistently outperform pure LLM approaches", 1),
])

# ── 22. QR 2023: First Evaluations ───────────────────────────────────────────
add_content_slide("QR 2023: First Systematic Evaluations", [
    "3 of 15 papers engaged directly with LLMs (12 papers: classical QR methods only)",
    "Paper 07 — Cohn (U Leeds): ChatGPT-4 on RCC-8 spatial reasoning",
    ("Composition table queries: 71.94% accuracy", 1),
    ("Preferred compositions (requires geometric judgment): 40.82% accuracy", 1),
    ("Conceptual neighborhood (recognition): ~90% accuracy", 1),
    ("Systematic failures on transitive composition chains — genuine spatial inference, not pattern matching", 1),
    "Paper 13 — Maxwell et al. (PARC/SRI): LLMs as physical system topology generator",
    ("ChatGPT proposes Modelica component connection topologies from NL descriptions", 1),
    ("'Sloppy but general': approximately correct far more often than random", 1),
    ("Qualitative simulation used as repair oracle: diagnoses missing components by behavior mismatch", 1),
    "Paper 14 — Forbus (Northwestern): Explicit rejection of LLMs as QR knowledge bases",
    ("'LLMs make poor knowledge bases: not grounded; optimizing for plausibility, not correctness'", 1),
    ("Proposes ontology-anchored QP theory as the alternative", 1),
])

# ── 23. Forbus Quote ──────────────────────────────────────────────────────────
add_quote_slide(
    "The Case Against LLMs as QR Knowledge Bases",
    "LLMs make poor knowledge bases for two reasons. First, their exposure to language is not grounded "
    "in the everyday world. Second, their success criterion is generating statistically plausible text, "
    "not correct reasoning. As the confabulation problems with LLMs show, these are at best only correlated.",
    "Kenneth Forbus, QR 2023 (Paper 14)"
)

# ── 24. QR 2024: LLMs in the Pipeline ─────────────────────────────────────────
add_content_slide("QR 2024: LLMs as Pipeline Components", [
    "2 of 9 papers engaged with LLMs/foundation models (7 papers: classical QR methods)",
    "Paper 3 — Suzuki & Yoshioka (Tohoku): GPT-4 as NL retriever for QP model construction",
    ("Pipeline: NL problem text → GPT-4 → physical objects + processes → QPT database lookup → QPT simulator", 1),
    ("LLM identifies what is involved; symbolic QPT does the reasoning", 1),
    ("Evaluated on introductory physics (heat flow, fluid dynamics, projectile motion)", 1),
    ("Key insight: LLMs have physical commonsense breadth; QPT has formal precision", 1),
    "Paper 9 — Keser et al. (Lübeck/Stockholm): CLIP analyzed for ontological commitments",
    ("Studies what ontological categories CLIP implicitly encodes in its embedding space", 1),
    ("Hierarchical clustering over structured image sets vs. WordNet, part-whole hierarchies", 1),
    ("CLIP reliably encodes basic-level categories; inconsistent on abstract properties and part-whole", 1),
    ("Opens research direction: QR ontological analysis tools applied to foundation models", 1),
])

# ── 25. QR 2025: Rapid Acceleration ──────────────────────────────────────────
add_content_slide("QR 2025: Rapid Acceleration", [
    "6 of 10 papers engage LLMs or foundation models — a dramatic shift from 2023",
    "Three roles identified:",
    ("Reasoning engines: LLMs evaluated on QR tasks directly (papers 0004, 0007, 0008)", 1),
    ("Translators/retrievers: LLMs as NL → formal structure front-ends (papers 0001, 0010)", 1),
    ("Modeled agents: LLMs as subjects of study in qualitative frameworks (paper 0005)", 1),
    "Large Reasoning Models (o1, o3) show dramatic improvements:",
    ("Chain-of-thought at inference time; not just pattern matching", 1),
    ("Best result: o1 achieves 0.92 accuracy on spatial QR benchmark (vs. GPT-4's 0.72)", 1),
    "New concern: epistemic alignment",
    ("Behavioral alignment (doing the right thing) may be insufficient", 1),
    ("We may also need agents that reason correctly about what they know and don't know", 1),
])

# ── 26. QR 2025: Key Papers ───────────────────────────────────────────────────
add_content_slide("QR 2025: Key Papers in Detail", [
    "Paper 0001 (Forbus, Northwestern): NL input + QP model library → hybrid QR system",
    ("Forbus acknowledges LLMs as useful NL → QPT translators while maintaining formal QP core", 1),
    "Paper 0004 (Ansah et al.): VLMs on mechanical comprehension (gear/pulley puzzles)",
    ("GPT-4V and similar models evaluated; most fail systematic compositional reasoning", 1),
    "Paper 0007 (Cohn & Blackwell): 28 LLMs benchmarked on spatial reasoning",
    ("Largest systematic comparison to date; LRMs (o1, o3) dominate", 1),
    ("Smaller models fail qualitatively differently from larger ones", 1),
    "Paper 0008 (Peng & Lai): Multi-agent QR with LLM supervision",
    ("LLM coordinates multi-agent qualitative simulation; qualitative feedback improves LLM decisions", 1),
    "Paper 0005 (Liza): Epistemic alignment for LLM-based AI systems",
    ("QR framework for evaluating whether LLM agents know what they don't know", 1),
    "Paper 0010 (Suzuki et al.): Extended LLM-QPT pipeline",
    ("Follow-up to QR2024 paper 3 — broader coverage, more evaluation", 1),
])

# ── 27. The Consistent Pattern ────────────────────────────────────────────────
add_two_column_slide(
    "The Consistent Pattern: Hybrid Neuro-Symbolic Architectures",
    "What LLMs do well in QR",
    [
        "Parsing natural language descriptions of physical systems",
        "Identifying relevant objects, processes, phenomena",
        "Generating plausible component topologies as starting points",
        "Recognizing spatial relation vocabulary (DC, EC, PO...)",
        "Broad commonsense coverage across domains",
        "Reducing knowledge engineering bottleneck",
    ],
    "What symbolic QR provides",
    [
        "Formal soundness: all real behaviors are captured",
        "Correct compositional spatial inference (RCC-8 tables)",
        "Causal explanation of every prediction",
        "Diagnosis of errors in LLM-generated structures",
        "Behavior-based repair of incorrect topologies",
        "Grounded, verifiable knowledge representations",
    ]
)

# ── 28. Benchmark Performance Table ───────────────────────────────────────────
add_content_slide("LLM Performance on QR Tasks: 2023–2025", [
    "Spatial reasoning (RCC-8 composition table) — Cohn et al.:",
    ("ChatGPT-4 (2023): 71.94%  |  GPT-4 preferred compositions: 40.82%", 1),
    ("o1 (2025): ~92%  |  o3 approaches formal QR solver level", 1),
    ("Conceptual neighborhood (easiest): ~90% even for GPT-4 (recognition task)", 1),
    "Physical system design — Maxwell et al. (QR 2023):",
    ("Low-pass filter topology: mostly correct  |  Power train: mostly incorrect without repair", 1),
    ("With qualitative repair oracle: acceptable designs produced reliably", 1),
    "Mechanical comprehension (gears, pulleys) — Ansah et al. (QR 2025):",
    ("Most VLMs fail systematic compositional reasoning even on simple puzzles", 1),
    "QP model construction — Suzuki & Yoshioka (QR 2024/2025):",
    ("LLM retrieval step: reliable for object/process identification", 1),
    ("QPT reasoning step: requires symbolic system — LLMs alone unreliable", 1),
    "Pattern: recognition tasks (nearest-neighbor, vocabulary) → LLMs strong",
    ("Compositional inference tasks (transitivity, chains) → LLMs weak; LRMs improving fast", 1),
])

# ── 29. The Grounding Critique Over Time ──────────────────────────────────────
add_content_slide("The Grounding Critique: From Hayes to Forbus", [
    "1977 — Hayes: computers lack commonsense grounding in the physical world",
    ("Solution proposed: formal axiomatization of naïve physics", 1),
    "1984–2000 — QR community builds formal grounded representations",
    ("QPT, QSIM, RCC-8: symbolic but well-specified formal semantics", 1),
    "2023 — Forbus extends the critique to LLMs:",
    ("'Exposure to language is not grounded in the everyday world'", 1),
    ("'Success criterion is statistical plausibility, not correct reasoning'", 1),
    "2024 — Keser et al. provide empirical evidence:",
    ("CLIP's embedding space partially but imperfectly aligns with formal ontologies", 1),
    ("Basic categories: good  |  Abstract properties: inconsistent  |  Part-whole: poor", 1),
    "2025 — Liza extends to epistemic alignment:",
    ("Not just 'does the agent behave correctly?' but 'does it know what it doesn't know?'", 1),
    "The grounding problem is 48 years old — LLMs reshape but do not resolve it",
])

# ── 30. Research Directions ────────────────────────────────────────────────────
add_content_slide("Open Research Directions at the QR/LLM Intersection", [
    "Automated model generation at scale",
    ("LLMs reduce the KB bottleneck — can we automate full QP model construction?", 1),
    ("Key challenge: LLM output must be verified against formal ontological constraints", 1),
    "Grounding foundation models in QR representations",
    ("Fine-tuning on QR-annotated corpora; embedding quantity-space structure", 1),
    ("Using QR ontologies to audit and repair LLM conceptual commitments", 1),
    "LRMs vs. symbolic QR: convergence or complementarity?",
    ("o1/o3 approach formal QR accuracy — do they still need symbolic backing?", 1),
    ("Answer likely domain-dependent: safety-critical systems require formal guarantees", 1),
    "Epistemic alignment for autonomous QR agents",
    ("Agents that correctly represent their own uncertainty over qualitative models", 1),
    "Multi-scale and multi-domain QR with LLM assistance",
    ("Ecological systems, social dynamics, autonomous vehicles — all benefit from QR breadth", 1),
])

# ── 31. QR + LLMs: Timeline ───────────────────────────────────────────────────
add_content_slide("Timeline: QR and LLMs", [
    "1977 — Hayes' Naïve Physics Manifesto",
    "1984 — Foundational QR papers (de Kleer, Forbus, Kuipers) in AI journal",
    "1986 — QSIM algorithm; 1987 — First QR Workshop",
    "1991 — Compositional modeling (Falkenhainer & Forbus)",
    "~1990–1995 — Commercial applications: photocopier, automotive FMEA",
    "1996 — NASA Deep Space One: first autonomous AI-controlled spacecraft",
    "2006 — QR community articulates 20-year grand challenge visions",
    "2010s — DynaLearn deployed in European schools; CogSketch in cognitive science",
    "2023 — First systematic evaluations of ChatGPT-4 on QR tasks (QR2023 Workshop)",
    "2024 — LLM-QPT hybrid pipeline demonstrated (Suzuki & Yoshioka)",
    "2025 — 28 LLMs benchmarked; o1 achieves ~92% on spatial QR; epistemic alignment emerges",
    "Future — LRMs approach formal QR accuracy; hybrid architectures mature; grounding unresolved",
])

# ── 32. Summary ───────────────────────────────────────────────────────────────
add_content_slide("Summary", [
    "Qualitative Reasoning is a 60+ year AI subfield with deep mathematical roots",
    ("Sign algebra (economics, 1962) → QPT + QSIM + Confluences (1984) → compositional modeling (1991)", 1),
    ("RCC-8 spatial reasoning, DynaLearn education tools, NASA spacecraft diagnosis", 1),
    "QR's core value: sound, explainable, causal reasoning under incomplete information",
    ("No training data required — knowledge encoded as model fragments", 1),
    ("Every conclusion traceable to specific model structure", 1),
    "LLMs are reshaping how QR is practiced — but not replacing it",
    ("LLMs: natural language → structured objects/processes (breadth)", 1),
    ("Symbolic QR: formal inference, soundness guarantees, causal explanation (precision)", 1),
    ("LRMs (o1, o3) rapidly close benchmark gap — formal guarantees remain symbolic", 1),
    "The grounding problem Hayes identified in 1977 remains the central challenge",
    ("LLMs are trained on language, not the physical world — same critique, new context", 1),
    "Hybrid neuro-symbolic architectures are the community's current consensus path forward",
])

# ══════════════════════════════════════════════════════════════════════════════
# REFERENCES
# ══════════════════════════════════════════════════════════════════════════════

add_section_divider("References", "Primary sources used in this presentation")

add_content_slide("References — History and Foundations", [
    "Lancaster, K. (1962). The solution of qualitative comparative static problems. Quarterly Journal of Economics.",
    "Hayes, P. (1977). In defence of logic. IJCAI-1977. Also: The naïve physics manifesto (1979, 1985).",
    "de Kleer, J. & Brown, J. (1984). A qualitative physics based on confluences. Artificial Intelligence, 24, 7–83.",
    "Forbus, K. (1984). Qualitative process theory. Artificial Intelligence, 24, 85–168.",
    "Kuipers, B. (1986). Qualitative simulation. Artificial Intelligence, 29, 289–338.",
    "Falkenhainer, B. & Forbus, K. (1991). Compositional modeling: finding the right model for the job. Artificial Intelligence, 51, 95–143.",
    "Mavrovouniotis, M. & Stephanopoulos, G. (1988). Formal order-of-magnitude reasoning in process engineering. Computers & Chemical Engineering, 12, 867–880.",
    "Travé-Massuyes, L. (1992). Qualitative reasoning: history, status, perspectives. Knowledge Engineering Review, 7(4).",
    "Price, C. et al. (MONET network). (2006). Qualitative futures. Knowledge Engineering Review, 21(4), 317–357.",
])

add_content_slide("References — Applications and Software", [
    "Williams, B. & Nayak, P. (1996). A model-based approach to reactive self-configuring systems. AAAI-1996. [NASA DS-1]",
    "Hamscher, W. (1991). Modeling digital circuits for troubleshooting. Artificial Intelligence, 51, 223–271. [Xerox FMEA]",
    "Bredeweg, B. et al. (2009). DynaLearn — an intelligent learning environment for learning conceptual knowledge. AI Magazine.",
    "Forbus, K. et al. (2011). CogSketch: Sketch understanding for cognitive science research and for education. Topics in Cognitive Science.",
    "Forbus, K. (2011). Qualitative reasoning. Chapter 35 in Computing Handbook, 3rd ed. CRC Press.",
    "Kuipers, B. (1994). Qualitative Reasoning: Modeling and Simulation with Incomplete Knowledge. MIT Press.",
    "Forbus, K. (2019). Qualitative Representations: How People Reason and Learn about the Continuous World. MIT Press.",
    "Kragten, M. et al. (2024). GarpN: Bridging qualitative and quantitative modelling. QR 2024 Workshop.",
])

add_content_slide("References — LLMs in QR (2023–2025)", [
    "Cohn, A.G. (2023). Evaluating ChatGPT on qualitative spatial reasoning. QR 2023 Workshop.",
    "Maxwell, J. et al. (2023). Preliminary experiments using LLMs for design. QR 2023 Workshop.",
    "Forbus, K. (2023). Building domain theories from language-grounded ontologies. QR 2023 Workshop.",
    "Suzuki, H. & Yoshioka, F. (2024). Automating qualitative model construction with LLMs and physical laws. QR 2024 Workshop.",
    "Keser, T., Wolter, D. & Bhatt, M. (2024). Extracting ontological commitment from foundation models. QR 2024 Workshop.",
    "Forbus, K. (2025). NL input with qualitative process model libraries. QR 2025 Workshop.",
    "Cohn, A.G. & Blackwell, A. (2025). Benchmarking 28 LLMs on qualitative spatial reasoning. QR 2025 Workshop.",
    "Ansah, J. et al. (2025). VLMs on mechanical comprehension tasks. QR 2025 Workshop.",
    "Peng, X. & Lai, T. (2025). Multi-agent qualitative supervision with LLMs. QR 2025 Workshop.",
    "Liza, F.F. (2025). Epistemic alignment for qualitative reasoning agents. QR 2025 Workshop.",
    "Suzuki, H. et al. (2025). Extended LLM-QPT pipeline for physics problems. QR 2025 Workshop.",
])

# ── Save ───────────────────────────────────────────────────────────────────────
output_path = "/Users/pisan/bitbucket/pisanuw/qr2025/papers/qr_presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
